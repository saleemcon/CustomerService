import io
import os
import streamlit as st
from PyPDF2 import PdfReader
from azure.ai.inference import ChatCompletionsClient
from azure.ai.inference.models import SystemMessage, UserMessage
from azure.core.credentials import AzureKeyCredential
import fitz
import pytesseract
from PIL import Image
import pptx
import pickle

# HuggingFace Embeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain.embeddings import HuggingFaceEmbeddings

# Azure AI Foundry Configuration
AZURE_ENDPOINT = "https://salee-m81ple1k-francecentral.services.ai.azure.com/models"
AZURE_MODEL_NAME = "DeepSeek-V3"
AZURE_API_KEY = "95boWs4DWgPr5Ar1xWh4Gnt2Lg9rR6uX2yBhi0vv9RrdeVnycYr7JQQJ99BCAC5T7U2XJ3w3AAAAACOGpy6K"

st.set_page_config(page_title="Customer Service AI", layout="centered")

icon = Image.open("assets/Customer_Service.png")
col1, col2 = st.columns([3, 10])  # Give col1 more width
with col1:
    st.image(icon, width=250)  # Reduce to a realistic width that fits




client = ChatCompletionsClient(
    endpoint=AZURE_ENDPOINT,
    credential=AzureKeyCredential(AZURE_API_KEY),
)

# Text extraction functions
def extract_text_from_pdf(file_path):
    extracted_text = []
    ocr_text = []

    # Read the file once into bytes
    with open(file_path, "rb") as f:
        file_bytes = f.read()

    # Attempt text-layer extraction using PyPDF2
    reader = PdfReader(io.BytesIO(file_bytes))
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            extracted_text.append(f"\n--- Page {i+1} [Text Layer] ---\n{text.strip()}")
        else:
            # OCR fallback using PyMuPDF (render image and OCR)
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            pix = doc.load_page(i).get_pixmap(dpi=300)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            ocr_result = pytesseract.image_to_string(img, lang="eng+ara")
            if ocr_result.strip():
                ocr_text.append(f"\n--- Page {i+1} [OCR Layer] ---\n{ocr_result.strip()}")

    return "\n".join(extracted_text + ocr_text).strip()

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

# Embedding creation
def create_vector_store_from_folder(folder_path):
    full_text = ""
    for filename in os.listdir(folder_path):
        full_path = os.path.join(folder_path, filename)
        if filename.lower().endswith(".pdf"):
            full_text += "\n" + extract_text_from_pdf(full_path)
        elif filename.lower().endswith(('.ppt', '.pptx')):
            full_text += "\n" + extract_text_from_pptx(full_path)

    if not full_text.strip():
        return None

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
    chunks = text_splitter.split_text(full_text)

    embed_model = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    return FAISS.from_texts(chunks, embedding=embed_model)

# Streamlit UI
st.markdown("<h1 style='color:black; text-align:center;'>AI Customer Service</h1>", unsafe_allow_html=True)
st.markdown("---")

# Folder selection
st.markdown("<h3 style='color: black; margin-bottom: 0.2rem;'>📂 Choose a product folder:</h3>", unsafe_allow_html=True)
category = st.selectbox("", ["TV", "OneDrive", "Cookers", "SDA"])
folder_paths = {
    "TV": "manuals/TV",
    "OneDrive": "manuals/OneDrive",
    "Cookers": "manuals/Cookers",
    "SDA": "manuals/SDA"
}
selected_folder = folder_paths[category]
cache_name = f".cache/faiss_folder_{category}.pkl"

# Load or create vector store
vector_store = None
if os.path.exists(cache_name):
    with open(cache_name, "rb") as f:
        vector_store = pickle.load(f)
    st.success("✅ Loaded cached embeddings.")
else:
    with st.spinner("🔍 Processing and embedding documents..."):
        vector_store = create_vector_store_from_folder(selected_folder)
        if vector_store:
            os.makedirs(".cache", exist_ok=True)
            with open(cache_name, "wb") as f:
                pickle.dump(vector_store, f)
            st.success("✅ Folder processed and cached.")
        else:
            st.error("❌ No supported documents found in the selected folder.")

# Prompt input
st.markdown("<h3 style='color: black;'>💬 Ask a question based on the documents:</h3>", unsafe_allow_html=True)
user_prompt = st.text_area("", placeholder="e.g., What is the warranty policy?", height=100)

# Generate answer

if st.button("🚀 Generate Response"):
    if user_prompt.strip() and vector_store:
        with st.spinner("🧠 Thinking..."):
            docs = vector_store.similarity_search(user_prompt, k=5)
            context = "\n\n".join([doc.page_content for doc in docs])

            messages = [
                SystemMessage(content="You are a helpful assistant that answers customer questions based on product documents."),
                UserMessage(content=f"Context:\n{context}"),
                UserMessage(content=f"Customer Question: {user_prompt}")
            ]

            response = client.complete(
                stream=True,
                messages=messages,
                max_tokens=2048,
                temperature=0.7,
                top_p=0.95,
                model=AZURE_MODEL_NAME
            )

            answer = ""
            for update in response:
                if update.choices:
                    answer += update.choices[0].delta.content or ""

            st.success("✅ AI Response:")
            st.markdown(f"<div style='background-color:#f1f1f1;padding:15px;border-radius:10px'>{answer}</div>", unsafe_allow_html=True)
    else:
        st.warning("⚠️ Please enter a question first.")

st.markdown("---")
st.markdown("<div style='text-align:center; color: gray;'>Made with ❤️ using Azure + HuggingFace + Streamlit</div>", unsafe_allow_html=True)
